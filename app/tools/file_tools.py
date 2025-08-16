"""
文件工具模块 - 提供安全的文件读写和目录操作功能

该模块提供了一系列工具函数，用于：
1. 安全地读取文本和二进制文件
2. 列出目录内容
3. 安全路径处理
4. 支持各种文件格式的解析

注意: 此模块已经与file_system_tools.py合并为统一的文件处理模块。
"""

import os
import sys
import json
import csv
import logging
import mimetypes
import pathlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, BinaryIO, TextIO, Tuple
import base64
from io import StringIO, BytesIO
import traceback
import chardet
from langgraph.config import get_stream_writer
import matplotlib.pyplot as plt
import numpy as np
import time
import datetime
import re
import tempfile
from uuid import uuid4

# 添加对Excel文件的支持
try:
    import pandas as pd
except ImportError:
    pd = None

# 添加对Word文档的支持
try:
    import docx
except ImportError:
    docx = None

# 添加对老式Word文档的支持
try:
    import win32com.client
    has_win32com = True
except ImportError:
    has_win32com = False

# 添加新的导入 - PPT支持
try:
    import pptx
    has_pptx = True
except ImportError:
    has_pptx = False

# 添加新的导入 - PDF支持
try:
    import PyPDF2
    has_pypdf = True
except ImportError:
    has_pypdf = False
    try:
        import pdfplumber
        has_pdfplumber = True
    except ImportError:
        has_pdfplumber = False

# 引入工具注册机制
from app.tools.registry import register_tool
from app.core.file_manager import file_manager

# 添加RAG相关导入
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    TextLoader, 
    CSVLoader, 
    PyPDFLoader, 
    Docx2txtLoader,
    UnstructuredExcelLoader
)
from app.utils.silicon_embeddings import SiliconFlowEmbeddings

# 配置日志
logger = logging.getLogger(__name__)

# 配置安全路径
DEFAULT_ENCODING = "utf-8"


# 支持的图片格式
SUPPORTED_IMAGE_FORMATS = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']

# 设置工作目录和安全目录
# 获取项目根目录
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
# 设置主工作目录为data文件夹
WORKING_DIR = os.path.join(PROJECT_ROOT, "data")

# 定义data下的子目录
DATA_SUBDIRS = [
    "uploads",
    "temp",
    "knowledge",
    "tmp",
    "tasks",
    "generated",
    "files"
]

# 构建安全目录列表，包含data目录及其所有子目录
SAFE_DIRS = [WORKING_DIR]
for subdir in DATA_SUBDIRS:
    SAFE_DIRS.append(os.path.join(WORKING_DIR, subdir))

# 确保目录存在
for dir_path in SAFE_DIRS:
    if not os.path.exists(dir_path):
        try:
            os.makedirs(dir_path, exist_ok=True)
            logger.info(f"已创建目录: {dir_path}")
        except Exception as e:
            logger.warning(f"无法创建目录 {dir_path}: {e}")

# 文件目录类型映射
FILE_DIR_MAPPING = {
    "upload": "uploads",     # 上传的文件
    "generated": "generated", # 生成的文件
    "temp": "temp",          # 临时文件
    "file": "files",         # 通用文件
    "knowledge": "knowledge", # 知识文件
    "task": "tasks",         # 任务相关文件
    "tmp": "tmp"             # 临时文件(另一种)
}

# 文件类型目录映射
FILE_TYPE_DIRS = {
    "image": "images",
    "text": "documents",
    "pdf": "documents",
    "spreadsheet": "spreadsheets", 
    "archive": "archives",
    "audio": "media/audio",
    "video": "media/video",
    "data": "data"
}

def is_path_safe(file_path: str) -> bool:
    """检查文件路径是否安全
    
    Args:
        file_path: 文件路径
        
    Returns:
        路径是否安全
    """
    # 首先标准化路径
    abs_path = os.path.abspath(file_path)
    
    # 检查路径是否在允许的目录中
    for allowed_path in SAFE_DIRS:
        allowed_abs_path = os.path.abspath(allowed_path)
        # 使用os.path.commonpath来更准确地检查路径是否为子目录
        try:
            # commonpath在Windows和Linux上都能正确处理路径比较
            if os.path.commonpath([allowed_abs_path, abs_path]) == allowed_abs_path:
                # 检查路径是否存在
                if not os.path.exists(abs_path):
                    parent_dir = os.path.dirname(abs_path)
                    # 如果父目录存在，则可能是要创建新文件，允许操作
                    if os.path.exists(parent_dir):
                        logger.info(f"路径不存在但父目录存在，可能是新文件: {abs_path}")
                        return True
                    else:
                        logger.warning(f"路径及其父目录不存在: {abs_path}")
                        return False
                return True
        except ValueError:
            # commonpath在比较不同驱动器上的路径时可能抛出ValueError
            continue
    
    logger.warning(f"路径不在允许的目录中: {abs_path}")
    return False

def normalize_path(file_path: str) -> str:
    """标准化文件路径，兼容Windows和Linux平台
    
    Args:
        file_path: 文件路径
        
    Returns:
        标准化后的路径
    """
    # 转换为绝对路径(如果是相对路径)
    if not os.path.isabs(file_path):
        # 如果是相对路径但没有明确以.或..开头，添加./前缀
        if not file_path.startswith('.') and not file_path.startswith('/'):
            file_path = f"./{file_path}"
        # 将相对路径转换为绝对路径，基于工作目录
        file_path = os.path.abspath(os.path.join(WORKING_DIR, file_path))
    
    # 解析路径，处理..和.等特殊路径元素
    normalized_path = os.path.normpath(file_path)
    
    # 确保路径分隔符与当前操作系统一致
    normalized_path = os.path.normpath(normalized_path)
    
    return normalized_path


def get_file_mime_type(file_path: str) -> str:
    """获取文件MIME类型
    
    Args:
        file_path: 文件路径
        
    Returns:
        MIME类型字符串
    """
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type or "application/octet-stream"

# @register_tool(category="file")
def read_image_file(file_path: str) -> str:
    """读取图片文件，返回基本信息和base64编码（如果可能）
    
    Args:
        file_path: 图片文件路径
        
    Returns:
        图片信息的字符串描述
    """
    # 标准化路径
    file_path = normalize_path(file_path)
    
    # 安全检查
    if not is_path_safe(file_path):
        return f"错误: 无法访问路径 {file_path}"
    
    # 检查文件大小
    try:
        file_size = os.path.getsize(file_path)
        if file_size > 10 * 1024 * 1024:  # 10MB
            return f"错误: 图片文件过大 ({file_size / (1024 * 1024):.2f} MB > 10 MB)"
    except Exception as e:
        return f"错误: 无法获取文件大小 - {str(e)}"
    
    # 获取MIME类型
    mime_type = get_file_mime_type(file_path)
    if not mime_type or not mime_type.startswith("image/"):
        return f"错误: 文件 {file_path} 不是有效的图片文件"
    
    try:
        # 尝试使用PIL读取图片信息
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                width, height = img.size
                format_name = img.format
                mode = img.mode
                
                result = {
                    "文件路径": file_path,
                    "图片类型": format_name,
                    "尺寸": f"{width}x{height}",
                    "颜色模式": mode,
                    "文件大小": f"{file_size / 1024:.2f} KB"
                }
                
                # 转换为字符串
                info_str = "图片信息:\n"
                for key, value in result.items():
                    info_str += f"- {key}: {value}\n"
                
                return info_str
        except ImportError:
            # 如果PIL不可用，返回基本信息
            return f"图片信息 (未安装PIL):\n- 文件路径: {file_path}\n- 文件类型: {mime_type}\n- 文件大小: {file_size / 1024:.2f} KB"
    
    except Exception as e:
        logger.error(f"读取图片文件时出错: {e}")
        return f"错误: 无法读取图片文件 - {str(e)}"

# @register_tool(category="file")
def list_available_files(session_id: Optional[str] = None, file_type: Optional[str] = None) -> str:
    """列出可用的文件
    
    Args:
        session_id: 可选的会话ID过滤
        file_type: 可选的文件类型过滤（如csv, xlsx, txt等）
        
    Returns:
        文件列表的字符串表示
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": "正在获取可用文件列表..."})
        
    try:
        # 获取文件列表
        files = file_manager.get_all_files(session_id=session_id, file_type=file_type)
        
        if not files:
            if writer:
                writer({"custom_step": "没有找到符合条件的文件"})
            return "没有找到符合条件的文件。"
        
        # 格式化文件信息
        result = "可用文件列表:\n"
        for i, file_info in enumerate(files, 1):
            file_id = file_info.get("file_id", "未知ID")
            file_name = file_info.get("file_name", "未知名称")
            file_type = file_info.get("file_type", "未知类型")
            file_size = file_info.get("size", 0)
            upload_time = file_info.get("upload_time", "未知时间")
            source = file_info.get("source", "未知来源")
            
            # 确定文件所在目录
            location_dir = "未知目录"
            if "file_path" in file_info:
                for dir_type, dir_name in FILE_TYPE_DIRS.items():
                    if dir_name in file_info["file_path"]:
                        location_dir = dir_name
                        break
            
            # 格式化文件大小
            if file_size < 1024:
                size_str = f"{file_size} 字节"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.2f} MB"
            
            file_entry = f"{i}. ID: {file_id}\n   名称: {file_name}\n   类型: {file_type}\n   大小: {size_str}\n   上传时间: {upload_time}\n   来源: {source}\n   位置: {location_dir}\n"
            result += file_entry
        
        if writer:
            writer({"custom_step": f"找到 {len(files)} 个文件"})
            
        return result
    except Exception as e:
        error_msg = f"列出文件时出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

# @register_tool(category="file")
def get_file_details(file_id: str) -> str:
    """获取文件名、文件类型、文件大小、上传时间、来源、存储位置、元数据信息的字符串表示
    
    Args:
        file_id: 文件ID
        
    Returns:
        文件名、文件类型、文件大小、上传时间、来源、存储位置、元数据信息的字符串表示
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": f"正在获取文件 {file_id} 的详细信息..."})
    
    try:
        # 获取文件信息
        file_info = file_manager.get_file_info(file_id)
        
        if not file_info:
            if writer:
                writer({"custom_step": f"找不到ID为 {file_id} 的文件"})
            return f"找不到ID为 {file_id} 的文件。"
        
        # 确定文件所在目录
        location_dir = "未知目录"
        if "file_path" in file_info:
            for dir_type, dir_name in FILE_TYPE_DIRS.items():
                if dir_name in file_info["file_path"]:
                    location_dir = dir_name
                    break
        
        # 格式化文件信息
        result = f"文件详细信息 (ID: {file_id}):\n"
        result += f"- 文件名: {file_info.get('file_name', '未知')}\n"
        result += f"- 文件类型: {file_info.get('file_type', '未知')}\n"
        result += f"- 内容类型: {file_info.get('content_type', '未知')}\n"
        result += f"- 存储位置: {location_dir}\n"
        
        # 格式化文件大小
        file_size = file_info.get("size", 0)
        if file_size < 1024:
            size_str = f"{file_size} 字节"
        elif file_size < 1024 * 1024:
            size_str = f"{file_size / 1024:.2f} KB"
        else:
            size_str = f"{file_size / (1024 * 1024):.2f} MB"
        result += f"- 文件大小: {size_str}\n"
        
        result += f"- 上传时间: {file_info.get('upload_time', '未知')}\n"
        result += f"- 来源: {file_info.get('source', '未知')}\n"
        result += f"- 文件路径: {file_info.get('file_path', '未知')}\n"
        
        # 添加元数据信息
        metadata = file_info.get("metadata", {})
        if metadata:
            result += "- 元数据:\n"
            for key, value in metadata.items():
                result += f"  - {key}: {value}\n"
        
        if writer:
            writer({"custom_step": "成功获取文件详细信息"})
            
        return result
    except Exception as e:
        error_msg = f"获取文件详细信息时出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

# @register_tool(category="file")
def search_files(query: str, session_id: Optional[str] = None) -> str:
    """搜索文件，返回文件名、文件类型、文件大小、上传时间、来源、存储位置、元数据信息的字符串表示
    
    Args:
        query: 搜索关键词
        session_id: 可选的会话ID过滤
        
    Returns:
        搜索结果的字符串表示
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": f"正在搜索文件: {query}..."})
    
    try:
        # 执行搜索
        files = file_manager.search_files(query, session_id=session_id)
        
        if not files:
            if writer:
                writer({"custom_step": f"没有找到匹配 '{query}' 的文件"})
            return f"没有找到匹配 '{query}' 的文件。"
        
        # 格式化搜索结果
        result = f"搜索 '{query}' 的结果:\n"
        for i, file_info in enumerate(files, 1):
            file_id = file_info.get("file_id", "未知ID")
            file_name = file_info.get("file_name", "未知名称")
            file_type = file_info.get("file_type", "未知类型")
            source = file_info.get("source", "未知来源")
            
            result += f"{i}. ID: {file_id}, 名称: {file_name}, 类型: {file_type}, 来源: {source}\n"
        
        if writer:
            writer({"custom_step": f"找到 {len(files)} 个匹配文件"})
            
        return result
    except Exception as e:
        error_msg = f"搜索文件时出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

@register_tool(category="file")
def preview_file_content(file_id: str, max_lines: int = 20) -> str:
    """预览xlsx格式文件内容，支持多种文件格式包括文本、CSV、Excel、PDF、Word和PowerPoint
    
    Args:
        file_id: 文件ID
        max_lines: 最大预览行数（或页数，取决于文件类型）
        
    Returns:
        文件内容预览的字符串表示
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": f"正在预览文件 {file_id} 的内容..."})
    
    try:
        # 获取文件信息
        file_info = file_manager.get_file_info(file_id)
        
        if not file_info:
            if writer:
                writer({"custom_step": f"找不到ID为 {file_id} 的文件"})
            return f"找不到ID为 {file_id} 的文件。"
        
        file_path = file_info.get("file_path")
        file_type = file_info.get("file_type", "").lower()
        file_name = file_info.get("file_name", "未知文件")
        
        # 处理不同类型的文件
        content = ""
        
        # CSV文件处理
        if file_type in ["csv"]:
            if writer:
                writer({"custom_step": "检测到CSV文件，使用pandas读取"})
            
            try:
                df = pd.read_csv(file_path)
                preview_df = df.head(max_lines)
                content = f"CSV文件内容预览 (前{min(max_lines, len(df))}行):\n"
                content += preview_df.to_string(index=True)
            except Exception as csv_err:
                if writer:
                    writer({"custom_step": f"使用pandas读取CSV失败: {str(csv_err)}"})
                
                # 回退到普通文本读取
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    lines = [line.strip() for line in f.readlines()[:max_lines]]
                    content = f"文件内容预览 (前{len(lines)}行):\n" + "\n".join(lines)
        
        # Excel文件处理
        elif file_type in ["xlsx", "xls"]:
            if writer:
                writer({"custom_step": "检测到Excel文件，使用pandas读取"})
            
            try:
                df = pd.read_excel(file_path)
                preview_df = df.head(max_lines)
                content = f"Excel文件内容预览 (前{min(max_lines, len(df))}行):\n"
                content += preview_df.to_string(index=True)
            except Exception as excel_err:
                if writer:
                    writer({"custom_step": f"使用pandas读取Excel失败: {str(excel_err)}"})
                content = f"无法预览Excel文件 {file_name}，文件可能损坏或格式不支持。"
        
        # DOCX文件处理
        elif file_type == "docx":
            if writer:
                writer({"custom_step": "检测到DOCX文件，使用python-docx读取"})
            
            if docx:
                try:
                    doc = docx.Document(file_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    # 限制段落数
                    preview_paragraphs = paragraphs[:max_lines]
                    content = f"DOCX文件内容预览 (前{len(preview_paragraphs)}段):\n\n"
                    content += "\n\n".join(preview_paragraphs)
                    
                    # 如果有表格，也预览表格内容
                    if doc.tables and len(content) < 5000:  # 限制内容长度
                        content += "\n\n表格内容预览:\n"
                        for i, table in enumerate(doc.tables[:3]):  # 最多预览3个表格
                            content += f"\n表格 {i+1}:\n"
                            for row in table.rows[:5]:  # 每个表格最多预览5行
                                content += " | ".join([cell.text for cell in row.cells]) + "\n"
                except Exception as docx_err:
                    if writer:
                        writer({"custom_step": f"读取DOCX文件失败: {str(docx_err)}"})
                    content = f"无法预览DOCX文件 {file_name}，文件可能损坏或格式不支持。"
            else:
                content = f"无法预览DOCX文件 {file_name}，因为python-docx库未安装。"
        
        # DOC文件处理（旧版Word格式）
        elif file_type == "doc":
            if writer:
                writer({"custom_step": "检测到DOC文件，尝试使用win32com或其他方法读取"})
            
            if has_win32com:
                try:
                    # 使用win32com读取doc文件
                    import win32com.client
                    import os
                    import tempfile
                    
                    # 创建临时目录保存转换后的文件
                    temp_dir = tempfile.mkdtemp()
                    temp_docx_path = os.path.join(temp_dir, "temp.docx")
                    
                    # 将doc转换为docx
                    word = win32com.client.Dispatch("Word.Application")
                    word.visible = False
                    doc = word.Documents.Open(os.path.abspath(file_path))
                    doc.SaveAs(os.path.abspath(temp_docx_path), 16)  # 16表示docx格式
                    doc.Close()
                    word.Quit()
                    
                    # 读取转换后的docx文件
                    docx_doc = docx.Document(temp_docx_path)
                    paragraphs = [p.text for p in docx_doc.paragraphs if p.text.strip()]
                    preview_paragraphs = paragraphs[:max_lines]
                    content = f"DOC文件内容预览 (前{len(preview_paragraphs)}段):\n\n"
                    content += "\n\n".join(preview_paragraphs)
                    
                    # 清理临时文件
                    try:
                        os.remove(temp_docx_path)
                        os.rmdir(temp_dir)
                    except:
                        pass
                except Exception as doc_err:
                    if writer:
                        writer({"custom_step": f"使用win32com读取DOC文件失败: {str(doc_err)}"})
                    # 尝试使用TextLoader作为备选方案
                    try:
                        from langchain_community.document_loaders import Docx2txtLoader
                        loader = Docx2txtLoader(file_path)
                        docs = loader.load()
                        text_content = "\n\n".join([doc.page_content for doc in docs])
                        preview_lines = text_content.split("\n")[:max_lines]
                        content = f"DOC文件内容预览 (使用备选方法, 前{len(preview_lines)}行):\n"
                        content += "\n".join(preview_lines)
                    except Exception as backup_err:
                        content = f"无法预览DOC文件 {file_name}，尝试的所有方法均失败。"
            else:
                try:
                    # 尝试使用docx2txt作为备选方案
                    from langchain_community.document_loaders import Docx2txtLoader
                    loader = Docx2txtLoader(file_path)
                    docs = loader.load()
                    text_content = "\n\n".join([doc.page_content for doc in docs])
                    preview_lines = text_content.split("\n")[:max_lines]
                    content = f"DOC文件内容预览 (前{len(preview_lines)}行):\n"
                    content += "\n".join(preview_lines)
                except Exception as alt_doc_err:
                    content = f"无法预览DOC文件 {file_name}，因为win32com未安装且备选方案也失败。{str(alt_doc_err)}"
        
        # PDF文件处理
        elif file_type == "pdf":
            if writer:
                writer({"custom_step": "检测到PDF文件，读取内容"})
            
            if has_pypdf:
                try:
                    # 使用PyPDF2读取PDF
                    with open(file_path, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        num_pages = len(pdf_reader.pages)
                        max_pages = min(num_pages, max_lines)
                        
                        content = f"PDF文件内容预览 ({file_name}, 共{num_pages}页):\n\n"
                        for i in range(max_pages):
                            page = pdf_reader.pages[i]
                            page_text = page.extract_text()
                            # 限制每页文本长度
                            if len(page_text) > 1000:
                                page_text = page_text[:997] + "..."
                            content += f"=== 第 {i+1} 页 ===\n{page_text}\n\n"
                except Exception as pdf_err:
                    if writer:
                        writer({"custom_step": f"使用PyPDF2读取PDF失败: {str(pdf_err)}"})
                    content = f"无法预览PDF文件 {file_name}，文件可能损坏或格式不支持。"
            elif has_pdfplumber:
                try:
                    # 使用pdfplumber作为备选方案
                    with pdfplumber.open(file_path) as pdf:
                        num_pages = len(pdf.pages)
                        max_pages = min(num_pages, max_lines)
                        
                        content = f"PDF文件内容预览 ({file_name}, 共{num_pages}页):\n\n"
                        for i in range(max_pages):
                            page = pdf.pages[i]
                            page_text = page.extract_text() or "[此页无文本内容]"
                            # 限制每页文本长度
                            if len(page_text) > 1000:
                                page_text = page_text[:997] + "..."
                            content += f"=== 第 {i+1} 页 ===\n{page_text}\n\n"
                except Exception as plumber_err:
                    if writer:
                        writer({"custom_step": f"使用pdfplumber读取PDF失败: {str(plumber_err)}"})
                    content = f"无法预览PDF文件 {file_name}，文件可能损坏或格式不支持。"
            else:
                try:
                    # 尝试使用langchain的PyPDFLoader
                    from langchain_community.document_loaders import PyPDFLoader
                    loader = PyPDFLoader(file_path)
                    docs = loader.load()
                    content = f"PDF文件内容预览 ({file_name}, 共{len(docs)}页):\n\n"
                    for i, doc in enumerate(docs[:max_lines]):
                        page_text = doc.page_content
                        if len(page_text) > 1000:
                            page_text = page_text[:997] + "..."
                        content += f"=== 第 {i+1} 页 ===\n{page_text}\n\n"
                except Exception as lc_pdf_err:
                    content = f"无法预览PDF文件 {file_name}，因为缺少PDF解析库。请安装PyPDF2或pdfplumber。"
        
        # PPT文件处理
        elif file_type in ["pptx"]:
            if writer:
                writer({"custom_step": "检测到PPTX文件，读取内容"})
            
            if has_pptx:
                try:
                    # 使用python-pptx读取PPTX
                    presentation = pptx.Presentation(file_path)
                    content = f"PPTX文件内容预览 ({file_name}, 共{len(presentation.slides)}张幻灯片):\n\n"
                    
                    # 限制幻灯片数量
                    max_slides = min(len(presentation.slides), max_lines)
                    for i, slide in enumerate(presentation.slides[:max_slides]):
                        content += f"=== 幻灯片 {i+1} ===\n"
                        
                        # 提取标题（如果有）
                        if slide.shapes.title:
                            content += f"标题: {slide.shapes.title.text}\n"
                        
                        # 提取文本框内容
                        slide_texts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_texts.append(shape.text)
                        
                        if slide_texts:
                            content += "内容:\n" + "\n".join(slide_texts) + "\n\n"
                        else:
                            content += "[此幻灯片没有可提取的文本内容]\n\n"
                except Exception as pptx_err:
                    if writer:
                        writer({"custom_step": f"读取PPTX文件失败: {str(pptx_err)}"})
                    content = f"无法预览PPTX文件 {file_name}，文件可能损坏或格式不支持。"
            else:
                content = f"无法预览PPTX文件 {file_name}，因为python-pptx库未安装。"
        
        # PPT文件处理（旧版PowerPoint格式）
        elif file_type == "ppt":
            if writer:
                writer({"custom_step": "检测到PPT文件，尝试使用备选方法"})
            
            if has_win32com:
                try:
                    # 使用win32com读取ppt文件并转换
                    import win32com.client
                    import os
                    import tempfile
                    
                    # 创建临时目录保存转换后的文件
                    temp_dir = tempfile.mkdtemp()
                    temp_pptx_path = os.path.join(temp_dir, "temp.pptx")
                    
                    # 将ppt转换为pptx
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    presentation = powerpoint.Presentations.Open(os.path.abspath(file_path))
                    presentation.SaveAs(os.path.abspath(temp_pptx_path))
                    presentation.Close()
                    powerpoint.Quit()
                    
                    # 使用python-pptx读取转换后的文件
                    if has_pptx:
                        presentation = pptx.Presentation(temp_pptx_path)
                        content = f"PPT文件内容预览 ({file_name}, 共{len(presentation.slides)}张幻灯片):\n\n"
                        
                        max_slides = min(len(presentation.slides), max_lines)
                        for i, slide in enumerate(presentation.slides[:max_slides]):
                            content += f"=== 幻灯片 {i+1} ===\n"
                            
                            if slide.shapes.title:
                                content += f"标题: {slide.shapes.title.text}\n"
                            
                            slide_texts = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    slide_texts.append(shape.text)
                            
                            if slide_texts:
                                content += "内容:\n" + "\n".join(slide_texts) + "\n\n"
                            else:
                                content += "[此幻灯片没有可提取的文本内容]\n\n"
                    else:
                        content = f"PPT文件已转换为PPTX，但无法预览，因为python-pptx库未安装。"
                    
                    # 清理临时文件
                    try:
                        os.remove(temp_pptx_path)
                        os.rmdir(temp_dir)
                    except:
                        pass
                except Exception as ppt_err:
                    if writer:
                        writer({"custom_step": f"处理PPT文件失败: {str(ppt_err)}"})
                    content = f"无法预览PPT文件 {file_name}，文件可能损坏或格式不支持。"
            else:
                content = f"无法预览PPT文件 {file_name}，因为缺少必要的库。请安装pywin32和python-pptx。"
        
        # 文本文件处理
        elif file_type in ["txt", "json", "xml", "md", "py", "js", "html", "css"]:
            if writer:
                writer({"custom_step": "检测到文本文件，直接读取内容"})
            
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    lines = [line.rstrip() for line in f.readlines()[:max_lines]]
                    content = f"文件内容预览 (前{len(lines)}行):\n" + "\n".join(lines)
            except Exception as txt_err:
                if writer:
                    writer({"custom_step": f"读取文本文件失败: {str(txt_err)}"})
                content = f"无法预览文本文件 {file_name}，文件可能损坏或格式不支持。"
        
        # 其他文件类型
        else:
            if writer:
                writer({"custom_step": f"不支持的文件类型预览: {file_type}"})
            content = f"不支持预览文件类型: {file_type}。文件: {file_name}"
        
        if writer:
            writer({"custom_step": "成功生成文件内容预览"})
            
        return content
    except Exception as e:
        error_msg = f"预览文件内容时出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

# @register_tool(category="file")
def get_current_session_files(session_id: str) -> str:
    """获取指定会话的文件列表
    
    Args:
        session_id: 会话ID
        
    Returns:
        会话文件列表的字符串表示
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": f"正在获取会话 {session_id} 的文件列表..."})
    
    try:
        # 直接使用传入的session_id
        if not session_id:
            if writer:
                writer({"custom_step": "未提供有效的会话ID"})
            return "错误: 未提供有效的会话ID"
        
        # 获取会话文件
        files = file_manager.get_session_files(session_id)
        
        if not files:
            if writer:
                writer({"custom_step": f"会话 {session_id} 没有关联的文件"})
            return f"会话 {session_id} 没有关联的文件。请先上传文件。"
        
        # 格式化文件列表
        result = f"会话 {session_id} 的文件列表:\n"
        for i, file_info in enumerate(files, 1):
            file_id = file_info.get("file_id", "未知ID")
            file_name = file_info.get("file_name", "未知名称")
            file_type = file_info.get("file_type", "未知类型")
            file_size = file_info.get("size", 0)
            upload_time = file_info.get("upload_time", "未知时间")
            
            # 格式化文件大小
            if file_size < 1024:
                size_str = f"{file_size} 字节"
            elif file_size < 1024 * 1024:
                size_str = f"{file_size / 1024:.2f} KB"
            else:
                size_str = f"{file_size / (1024 * 1024):.2f} MB"
            
            file_entry = f"{i}. ID: {file_id}\n   名称: {file_name}\n   类型: {file_type}\n   大小: {size_str}\n   上传时间: {upload_time}\n"
            result += file_entry
        
        # 添加使用说明
        result += "\n您可以通过文件ID使用其他文件系统工具进行操作，例如:\n"
        result += "- 使用preview_file_content工具预览文件内容\n"
        result += "- 使用get_file_details工具获取文件详细信息\n"
        
        if writer:
            writer({"custom_step": f"找到 {len(files)} 个文件"})
            
        return result
    except Exception as e:
        error_msg = f"获取会话文件时出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

@register_tool(category="file")
def search_documents_rag(
    query: str, 
    file_ids: str, 
    chunk_size: int = 600, 
    chunk_overlap: int = 100, 
    top_k: int = 10,
    rerank_with_llm: bool = True,
    reranked_n: int = 5
) -> str:
    """轻量级docx文档检索工具，在指定文档中查询检索地质背景相关内容请使用此工具
    
    Args:
        query: 查询内容
        file_ids: 要检索的文件ID列表，多个ID用逗号分隔
        chunk_size: 文档分块大小，默认600
        chunk_overlap: 文档分块重叠大小，默认100
        top_k: 返回的最相关文档数量，默认10
        rerank_with_llm: 是否使用LLM对结果进行重排序，默认False
        reranked_n: LLM重排序后保留的结果数量，默认5
        
    Returns:
        检索结果字符串，包含相关文本片段和元数据
    """
    # 获取流写入器，用于显示处理步骤
    writer = get_stream_writer()
    
    if writer:
        writer({"custom_step": f"开始执行RAG检索，查询内容: {query[:50]}..."})
        
    try:
        # 解析文件ID列表
        file_id_list = [id.strip() for id in file_ids.split(",") if id.strip()]
        
        if not file_id_list:
            if writer:
                writer({"custom_step": "未提供有效的文件ID"})
            return "错误: 未提供有效的文件ID"
        
        if writer:
            writer({"custom_step": f"将处理 {len(file_id_list)} 个文件"})
            
        # 加载所有文档
        documents = []
        for file_id in file_id_list:
            try:
                # 获取文件信息
                file_info = file_manager.get_file_info(file_id)
                if not file_info:
                    if writer:
                        writer({"custom_step": f"找不到ID为 {file_id} 的文件"})
                    continue
                
                file_path = file_info.get("file_path")
                file_name = file_info.get("file_name", "未知文件名")
                file_type = file_info.get("file_type", "").lower()
                
                if writer:
                    writer({"custom_step": f"处理文件: {file_name} (类型: {file_type})"})
                
                # 根据文件类型加载文档
                file_docs = []
                
                if file_type in ["txt", "md", "py", "js", "html", "css", "json", "xml"]:
                    # 文本文件处理
                    loader = TextLoader(file_path, encoding="utf-8", autodetect_encoding=True)
                    file_docs = loader.load()
                elif file_type == "pdf":
                    # PDF文件处理
                    loader = PyPDFLoader(file_path)
                    file_docs = loader.load()
                elif file_type in ["docx", "doc"]:
                    # Word文件处理
                    loader = Docx2txtLoader(file_path)
                    file_docs = loader.load()
                elif file_type in ["csv"]:
                    # CSV文件处理
                    loader = CSVLoader(file_path)
                    file_docs = loader.load()
                elif file_type in ["xlsx", "xls"]:
                    # Excel文件处理
                    loader = UnstructuredExcelLoader(file_path)
                    file_docs = loader.load()
                else:
                    if writer:
                        writer({"custom_step": f"不支持的文件类型: {file_type}"})
                    continue
                
                # 添加文件元数据
                for doc in file_docs:
                    if not hasattr(doc, "metadata"):
                        doc.metadata = {}
                    doc.metadata["file_id"] = file_id
                    doc.metadata["file_name"] = file_name
                    doc.metadata["file_type"] = file_type
                
                documents.extend(file_docs)
                
                if writer:
                    writer({"custom_step": f"成功加载文件: {file_name}，获取 {len(file_docs)} 个文档段"})
                
            except Exception as e:
                if writer:
                    writer({"custom_step": f"处理文件 {file_id} 时出错: {str(e)}"})
                logger.error(f"处理文件 {file_id} 时出错: {str(e)}")
                continue
        
        if not documents:
            if writer:
                writer({"custom_step": "未能加载任何有效文档"})
            return "未能加载任何有效文档，请检查文件ID和文件格式是否支持。"
        
        # 文档切片
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            is_separator_regex=False,
        )
        
        if writer:
            writer({"custom_step": f"使用参数 chunk_size={chunk_size}, chunk_overlap={chunk_overlap} 进行文档切片"})
        
        chunks = text_splitter.split_documents(documents)
        
        if writer:
            writer({"custom_step": f"文档切片完成，共生成 {len(chunks)} 个文本块"})
        
        # 创建嵌入模型实例
        try:
            embeddings = SiliconFlowEmbeddings()
            if writer:
                writer({"custom_step": "成功初始化SiliconFlowEmbeddings"})
        except Exception as e:
            logger.error(f"初始化SiliconFlowEmbeddings失败: {str(e)}")
            if writer:
                writer({"custom_step": f"初始化嵌入模型失败: {str(e)}"})
            return f"初始化嵌入模型失败: {str(e)}"
        
        # 创建向量存储
        try:
            if writer:
                writer({"custom_step": "正在创建向量存储..."})
            
            vectorstore = FAISS.from_documents(chunks, embeddings)
            
            if writer:
                writer({"custom_step": "向量存储创建完成"})
        except Exception as e:
            logger.error(f"创建向量存储失败: {str(e)}")
            if writer:
                writer({"custom_step": f"创建向量存储失败: {str(e)}"})
            return f"创建向量存储失败: {str(e)}"
        
        # 执行查询
        try:
            if writer:
                writer({"custom_step": f"执行查询: {query}"})
            
            search_results = vectorstore.similarity_search_with_score(query, k=top_k)
            
            if writer:
                writer({"custom_step": f"查询完成，找到 {len(search_results)} 个相关结果"})
        except Exception as e:
            logger.error(f"执行查询失败: {str(e)}")
            if writer:
                writer({"custom_step": f"执行查询失败: {str(e)}"})
            return f"执行查询失败: {str(e)}"
            
        # 使用LLM进行结果重排序
        if rerank_with_llm and len(search_results) > 0:
            if writer:
                writer({"custom_step": "正在使用LLM对结果进行重排序..."})
                
            try:
                # 导入LLM
                from app.utils.qwen_chat import SFChatOpenAI
                from app.core.config import ConfigManager
                # 初始化LLM
                config_manager = ConfigManager()
                model_config = config_manager.get_model_config()
                llm = SFChatOpenAI(
                    model=model_config.get("model_name", "Qwen/Qwen2.5-72B-Instruct-128K"), 
                    temperature=model_config.get("temperature", 0.1)
                )
                
                # 准备重排序数据
                rerank_items = []
                for i, (doc, score) in enumerate(search_results):
                    content = doc.page_content
                    file_name = doc.metadata.get("file_name", "未知文件")
                    page_number = doc.metadata.get("page", "")
                    
                    rerank_items.append({
                        "index": i,
                        "content": content,
                        "file_name": file_name,
                        "page_info": page_number,
                        "vector_score": score
                    })
                
                # 构建LLM提示
                rerank_prompt = f"""你是一个专业的文档相关性评估专家。我将给你一个查询和一些文档片段，请你评估每个片段与查询的相关性，并按照相关性从高到低的顺序对它们进行排序。
                
查询: {query}

文档片段:
"""
                for i, item in enumerate(rerank_items):
                    rerank_prompt += f"\n文档{i+1}: {item['content']}\n"
                
                rerank_prompt += """\n请对每个文档片段与查询的相关性进行评分，分数范围为0-10，0表示完全不相关，10表示非常相关。然后按照分数从高到低的顺序重新排列文档，输出格式如下:

文档评分:
文档1: 8
文档5: 7
文档3: 6
...
(请对所有文档评分并排序)
"""
                
                # 调用LLM进行评分
                if writer:
                    writer({"custom_step": "正在使用LLM评估文档相关性..."})
                
                rerank_response = llm.invoke(rerank_prompt)
                rerank_text = rerank_response.content
                
                if writer:
                    writer({"custom_step": "LLM评估完成，正在解析结果..."})
                
                # 解析LLM返回的评分和排序
                reranked_indices = []
                try:
                    # 提取评分部分
                    if "文档评分" in rerank_text:
                        scores_text = rerank_text.split("文档评分:")[1].strip()
                    else:
                        scores_text = rerank_text
                    
                    # 使用正则表达式提取文档编号和评分
                    import re
                    scores_pattern = r'文档(\d+):\s*(\d+\.?\d*)'
                    scores_matches = re.findall(scores_pattern, scores_text)
                    
                    # 创建评分列表
                    doc_scores = []
                    for doc_num, score in scores_matches:
                        doc_idx = int(doc_num) - 1  # 转换为0-索引
                        if 0 <= doc_idx < len(rerank_items):
                            doc_scores.append((doc_idx, float(score)))
                    
                    # 按评分排序
                    doc_scores.sort(key=lambda x: x[1], reverse=True)
                    
                    # 提取排序后的索引
                    reranked_indices = [idx for idx, _ in doc_scores]
                    
                    if writer:
                        writer({"custom_step": f"LLM重排序完成，得到{len(reranked_indices)}个结果"})
                except Exception as parse_err:
                    logger.error(f"解析LLM重排序结果失败: {str(parse_err)}")
                    if writer:
                        writer({"custom_step": f"解析重排序结果失败: {str(parse_err)}，将使用原始排序"})
                    reranked_indices = list(range(len(search_results)))
                
                # 限制结果数量
                reranked_indices = reranked_indices[:reranked_n]
                
                # 重排序结果
                reranked_results = [search_results[idx] for idx in reranked_indices if idx < len(search_results)]
                
                # 使用重排序后的结果
                search_results = reranked_results
                
                if writer:
                    writer({"custom_step": f"保留重排序后的前{len(search_results)}个结果"})
            except Exception as rerank_err:
                logger.error(f"LLM重排序失败: {str(rerank_err)}")
                if writer:
                    writer({"custom_step": f"LLM重排序过程中出错: {str(rerank_err)}，将使用原始排序"})
                # 保持原始排序
        
        # 格式化结果
        results_text = f"查询: {query}\n\n"
        
        if rerank_with_llm:
            results_text += f"在 {len(file_id_list)} 个文件中找到 {len(search_results)} 个经过LLM重排序的相关结果:\n\n"
        else:
            results_text += f"在 {len(file_id_list)} 个文件中找到 {len(search_results)} 个相关结果:\n\n"
        
        for i, (doc, score) in enumerate(search_results, 1):
            # 获取文档元数据
            file_name = doc.metadata.get("file_name", "未知文件")
            page_number = doc.metadata.get("page", "")
            page_info = f"页码: {page_number}" if page_number else ""
            
            # 格式化相似度分数
            similarity = 1.0 - float(score)  # 转换距离为相似度
            similarity_percentage = f"{similarity:.2%}"
            
            # 限制内容长度
            content = doc.page_content
            if len(content) > 500:
                content = content[:497] + "..."
            
            # 添加到结果
            results_text += f"结果 {i} (相关度: {similarity_percentage}):\n"
            results_text += f"- 来源: {file_name} {page_info}\n"
            results_text += f"- 内容:\n{content}\n\n"
        
        if writer:
            writer({"custom_step": "RAG检索完成，返回结果"})
            
        return results_text
    
    except Exception as e:
        error_msg = f"RAG搜索过程中出错: {str(e)}"
        logger.error(error_msg)
        if writer:
            writer({"custom_step": error_msg})
        return error_msg

# @register_tool(category="file")
def generate_test_image(title="测试图表", x_label="X轴", y_label="Y轴") -> str:
    """生成一个测试图表并保存
    
    Args:
        title: 图表标题
        x_label: X轴标签
        y_label: Y轴标签
        
    Returns:
        生成图表的路径和描述
    """
    # 获取流写入器
    writer = get_stream_writer()
    
    # 确保输出目录存在
    output_dir = os.path.join(WORKING_DIR, "generated")
    os.makedirs(output_dir, exist_ok=True)
    
    # 创建一个简单的图表
    plt.figure(figsize=(10, 6))
    
    # 发送开始生成图表的消息
    if writer:
        writer({"custom_step": f"开始生成图表: {title}"})
    
    # 生成数据
    x = np.linspace(0, 10, 100)
    y1 = np.sin(x)
    y2 = np.cos(x)
    
    # 绘制图表
    plt.plot(x, y1, 'b-', label='Sin(x)')
    plt.plot(x, y2, 'r--', label='Cos(x)')
    
    # 添加标题和标签
    plt.title(title)
    plt.xlabel(x_label)
    plt.ylabel(y_label)
    
    # 添加网格和图例
    plt.grid(True, alpha=0.3)
    plt.legend()
    
    # 发送绘制完成的消息
    if writer:
        writer({"custom_step": f"图表绘制完成，准备保存: {title}"})
    
    # 保存图表
    timestamp = int(time.time())
    output_filename = f"test_plot_{timestamp}.png"
    output_path = os.path.join(output_dir, output_filename)
    
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    
    # 注册文件到文件管理器
    from app.core.file_manager import file_manager
    file_info = file_manager.register_file(
        file_path=output_path,
        file_name=output_filename,
        file_type="png",
        source="generated",
        metadata={"description": f"由generate_test_image工具生成的测试图表: {title}"}
    )
    
    # 发送图片消息
    if writer:
        # 发送图片消息
        image_message = {
            "image_path": output_path,
            "title": f"生成的图表: {title}"
        }
        writer({"image_message": image_message})
        
        # 另外发送一个确认消息
        writer({"custom_step": f"图表已生成并保存: {output_filename}"})
    
    # 返回结果
    return f"已生成测试图表'{title}'并保存到文件: {output_path}\n文件ID: {file_info.get('file_id', '未知')}"

# 初始化模块
def init_module():
    """初始化文件工具模块"""
    logger.info(f"初始化文件工具模块，工作目录: {WORKING_DIR}")
    
    # 首先确保工作目录存在
    if not os.path.exists(WORKING_DIR):
        try:
            logger.info(f"工作目录不存在，正在创建: {WORKING_DIR}")
            os.makedirs(WORKING_DIR, exist_ok=True)
        except Exception as e:
            logger.error(f"创建工作目录失败: {str(e)}")
            raise RuntimeError(f"无法创建工作目录 {WORKING_DIR}: {str(e)}")
    
    # 确保所有安全目录存在
    for path in SAFE_DIRS:
        try:
            os.makedirs(path, exist_ok=True)
            logger.info(f"确保目录存在: {path}")
        except Exception as e:
            logger.error(f"创建目录 {path} 失败: {str(e)}")
    
    # 检查文件目录映射中的目录是否存在
    for dir_type, dir_name in FILE_DIR_MAPPING.items():
        dir_path = os.path.join(WORKING_DIR, dir_name)
        try:
            os.makedirs(dir_path, exist_ok=True)
            logger.info(f"确保文件类型目录存在: {dir_path}")
        except Exception as e:
            logger.error(f"创建文件类型目录 {dir_path} 失败: {str(e)}")
    
    # 检查文件类型目录是否存在
    for type_dir in FILE_TYPE_DIRS.values():
        # 对每个文件类型目录，确保它在主目录中存在
        for main_dir in [os.path.join(WORKING_DIR, dir_name) for dir_name in FILE_DIR_MAPPING.values()]:
            type_path = os.path.join(main_dir, type_dir)
            try:
                os.makedirs(type_path, exist_ok=True)
                logger.debug(f"确保文件类型子目录存在: {type_path}")
            except Exception as e:
                logger.warning(f"创建文件类型子目录 {type_path} 失败: {str(e)}")
    
    # 注册MIME类型
    mimetypes.init()
    
    logger.info(f"文件工具模块初始化完成，工作目录: {WORKING_DIR}")
    logger.info(f"允许的安全路径: {SAFE_DIRS}")

# 模块加载时自动初始化
init_module() 